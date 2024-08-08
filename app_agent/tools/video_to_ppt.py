import os
from pptx import Presentation
from pptx.util import Inches, Pt
from moviepy.editor import VideoFileClip
from PIL import Image


def get_video_files(folder_path):
    video_extensions = ['.mp4', '.avi', '.mov', '.wmv', '.mkv']
    video_files = [f for f in os.listdir(folder_path) if os.path.splitext(f)[1].lower() in video_extensions]
    return video_files


def save_first_frame_as_image(video_path, output_image_path):
    video_clip = VideoFileClip(video_path)
    first_frame = video_clip.get_frame(0)  # Get the first frame
    image = Image.fromarray(first_frame)
    image.save(output_image_path)
    video_clip.reader.close()
    if video_clip.audio:
        video_clip.audio.reader.close_proc()


def create_ppt_with_videos(folder_path, output_ppt_path):
    video_files = get_video_files(folder_path)

    # Create a presentation object
    prs = Presentation()
    slide_width, slide_height = prs.slide_width, prs.slide_height

    for video_file in video_files:
        title = os.path.splitext(video_file)[0]
        # Add a slide
        slide_layout = prs.slide_layouts[5]  # Use a blank slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Add video
        video_path = os.path.join(folder_path, video_file)
        first_frame_path = os.path.join(folder_path, f"{title}_first_frame.png")

        # Save the first frame of the video as an image
        save_first_frame_as_image(video_path, first_frame_path)

        # Insert the video with the first frame as the cover
        video_clip = VideoFileClip(video_path)

        # Calculate the size and position of the video to maintain aspect ratio
        video_width, video_height = video_clip.size
        aspect_ratio = video_width / video_height
        ppt_aspect_ratio = slide_width / slide_height

        # Set title position and size
        title_top = Inches(0.5)
        title_left = Inches(0.5)
        title_width = slide_width - Inches(1)
        title_height = Inches(1)

        # Set the title
        title_shape = slide.shapes.title
        if not title_shape:
            title_shape = slide.shapes.add_textbox(title_left, title_top, title_width, title_height)
        title_frame = title_shape.text_frame
        title_frame.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(24)  # Adjust the font size as needed

        # Calculate video position below the title with some margin
        video_top_margin = title_top + title_height   # Reduced the top margin between title and video
        if aspect_ratio > ppt_aspect_ratio:
            video_display_width = slide_width - Inches(1)  # Add some margin on sides
            video_display_height = (slide_width - Inches(1)) / aspect_ratio
        else:
            video_display_height = slide_height - video_top_margin - Inches(0.5)  # Reduced the bottom margin
            video_display_width = (slide_height - video_top_margin - Inches(0.5)) * aspect_ratio

        # Center the video horizontally and position it below the title
        left = (slide_width - video_display_width) / 2
        top = video_top_margin

        slide.shapes.add_movie(
            video_path, left, top, video_display_width, video_display_height, poster_frame_image=first_frame_path
        )

        # Close the video clip to release resources
        video_clip.reader.close()
        if video_clip.audio:
            video_clip.audio.reader.close_proc()

        # 删除封面图片
        if os.path.exists(first_frame_path):
            os.remove(first_frame_path)

    # Save the presentation
    prs.save(output_ppt_path)


# 使用示例
folder_path = 'C:/Users/86150/Videos/pptvideos'
output_ppt_path = 'C:/Users/86150/Videos/pptvideos/output.pptx'
create_ppt_with_videos(folder_path, output_ppt_path)
